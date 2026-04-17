"""
Генерация пояснительной записки и презентации по проекту «Фотоассистент».
Шаблон презентации: presentation-template.pptx (путь задаётся ниже).
Запуск из корня проекта:
  python scripts/generate_project_docs.py
Требует: pip install python-docx python-pptx
"""

from __future__ import annotations

import json
import os
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
# Шаблон: положите presentation-template.pptx в Downloads или укажите путь в PHOTOASSISTANT_PPTX_TEMPLATE
_default_pptx = Path.home() / "Downloads" / "presentation-template.pptx"
TEMPLATE_PPTX = Path(os.environ.get("PHOTOASSISTANT_PPTX_TEMPLATE", str(_default_pptx)))
OUT_PPTX = ROOT / "Презентация_Фотоассистент_v2.pptx"
OUT_DOCX = ROOT / "Пояснительная_записка_Фотоассистент.docx"

PPT_FONT = os.environ.get("PHOTOASSISTANT_PPTX_FONT", "Calibri")


def _apply_shape_font(shape, *, font_name: str = PPT_FONT, font_size: int | None = None) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            run.font.name = font_name
            if font_size is not None:
                run.font.size = PptPt(font_size)


def set_shape_text(shape, text: str, *, font_size: int | None = None) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.clear()
    if not tf.paragraphs:
        tf.add_paragraph()
    tf.paragraphs[0].text = text
    _apply_shape_font(shape, font_size=font_size)


def _resolve_pptx_template() -> Path:
    if TEMPLATE_PPTX.is_file():
        return TEMPLATE_PPTX
    local = ROOT / "presentation-template.pptx"
    if local.is_file():
        return local
    raise FileNotFoundError(
        f"Нет шаблона презентации. Ожидалось: {TEMPLATE_PPTX} или {local}. "
        "Скопируйте presentation-template.pptx в Downloads или в корень проекта, либо задайте PHOTOASSISTANT_PPTX_TEMPLATE."
    )


def _count_files(path: Path) -> int:
    if not path.is_dir():
        return 0
    return sum(1 for p in path.iterdir() if p.is_file())


def _load_model_metrics() -> dict:
    metrics_path = ROOT / "assets" / "model_metrics.json"
    if not metrics_path.is_file():
        return {}
    try:
        return json.loads(metrics_path.read_text(encoding="utf-8"))
    except Exception:
        return {}


def build_presentation() -> None:
    tpl = _resolve_pptx_template()
    pr = Presentation(str(tpl))
    slides = pr.slides
    metrics = _load_model_metrics()

    train_open = _count_files(ROOT / "training" / "data" / "prepared" / "train" / "open")
    train_closed = _count_files(ROOT / "training" / "data" / "prepared" / "train" / "closed")
    val_open = _count_files(ROOT / "training" / "data" / "prepared" / "val" / "open")
    val_closed = _count_files(ROOT / "training" / "data" / "prepared" / "val" / "closed")
    train_total = train_open + train_closed
    val_total = val_open + val_closed
    all_total = train_total + val_total
    train_pct = (train_total / all_total * 100.0) if all_total else 0.0
    val_pct = (val_total / all_total * 100.0) if all_total else 0.0
    open_total = train_open + val_open
    closed_total = train_closed + val_closed
    open_pct = (open_total / all_total * 100.0) if all_total else 0.0
    closed_pct = (closed_total / all_total * 100.0) if all_total else 0.0

    # Слайд 1 — титул
    s = slides[0]
    set_shape_text(s.shapes[1], "ФОТОАССИСТЕНТ ФОТОГРАФА", font_size=36)
    set_shape_text(
        s.shapes[3],
        "Десктоп-приложение на Python для автоматической сортировки снимков по экспозиции, резкости и закрытым глазам",
    )
    set_shape_text(
        s.shapes[4],
        "Выполнили:\nВолков Евгений\nНужина Арина\nМагнитогорск, 2026 г.",
    )

    # Слайд 2 — содержание (по требованию)
    s = slides[1]
    set_shape_text(s.shapes[1], "Содержание", font_size=34)
    set_shape_text(s.shapes[3], "1) Датасеты и подготовка", font_size=22)
    set_shape_text(
        s.shapes[5],
        "2) Модели: экспозиция, резкость, MediaPipe (глаза)\n"
        "3) ResNet-18 для глаз: обучение и графики\n"
        "4) Метрики качества\n"
        "5) Сравнение режимов и результаты\n"
        "6) Демонстрация приложения",
        font_size=18,
    )

    # Слайд 3 — актуальность (перенесено)
    s = slides[2]
    set_shape_text(s.shapes[1], "Актуальность проекта", font_size=34)
    set_shape_text(s.shapes[9], "Рост объёма цифровых фотографий", font_size=20)
    set_shape_text(
        s.shapes[10],
        "Съёмка сериями приводит к сотням и тысячам кадров; ручной отбор по качеству занимает много времени.",
        font_size=16,
    )
    set_shape_text(s.shapes[14], "Задача отбора портретов", font_size=20)
    set_shape_text(
        s.shapes[11],
        "Дополнительно важно отсеивать кадры с закрытыми глазами — это сложно и утомительно при ручном просмотре.",
        font_size=16,
    )

    # Слайд 4 — датасеты (по требованию)
    s = slides[3]
    set_shape_text(s.shapes[1], "Датасеты и подготовка", font_size=34)
    set_shape_text(s.shapes[4], "MRL Eye Dataset", font_size=22)
    set_shape_text(
        s.shapes[5],
        "Назначение: обучение ResNet-18 для классификации глаза (open/closed).\n"
        "Источник: https://mrl.cs.vsb.cz/eyedataset\n"
        "Подготовка: scripts/prepare_mrl_dataset.py → train/open|closed, val/open|closed",
        font_size=16,
    )
    set_shape_text(s.shapes[8], "Фото для анализа", font_size=22)
    set_shape_text(
        s.shapes[9],
        "Реальные .jpg/.jpeg/.png пользователя.\n"
        "Распределение в папки: Удачные / Неудачные/<причина>.",
        font_size=16,
    )
    set_shape_text(s.shapes[12], "«База данных» и проценты", font_size=22)
    set_shape_text(
        s.shapes[13],
        f"Train/Val: {train_total}/{val_total} ({train_pct:.1f}%/{val_pct:.1f}%).\n"
        f"Open/Closed: {open_total}/{closed_total} ({open_pct:.1f}%/{closed_pct:.1f}%).",
        font_size=16,
    )

    # Слайд 5 — модель 1: экспозиция
    s = slides[4]
    set_shape_text(s.shapes[1], "Модель: экспозиция", font_size=34)
    set_shape_text(s.shapes[4], "Метод", font_size=20)
    set_shape_text(s.shapes[5], "Пороговая оценка пересвета/недосвета по распределению яркости.", font_size=16)
    set_shape_text(s.shapes[8], "Метрики/параметры", font_size=20)
    set_shape_text(
        s.shapes[9],
        "exposure_low_thresh / exposure_high_thresh\nexposure_extreme_fraction — доля «выбитых» пикселей.\n"
        "Результат: OK / недосвет / пересвет.",
        font_size=16,
    )
    set_shape_text(s.shapes[12], "Где в коде", font_size=20)
    set_shape_text(s.shapes[13], "src/analysis/exposure.py + настройки в app_settings.json", font_size=16)
    set_shape_text(s.shapes[16], "Плюсы", font_size=20)
    set_shape_text(s.shapes[17], "Быстро, объяснимо, не требует обучения.", font_size=16)

    # Слайд 6 — модель 2: резкость
    s = slides[5]
    set_shape_text(s.shapes[1], "Модель: резкость", font_size=34)
    set_shape_text(s.shapes[2], "Метод", font_size=20)
    set_shape_text(s.shapes[3], "Дисперсия Лапласа (Laplacian variance) — оценка смаза/размытия.", font_size=16)
    set_shape_text(s.shapes[4], "Метрики/параметры", font_size=20)
    set_shape_text(s.shapes[5], "sharpness_threshold — порог; ниже → «смаз/размытие».", font_size=16)
    set_shape_text(s.shapes[6], "Где в коде", font_size=20)
    set_shape_text(s.shapes[7], "src/analysis/sharpness.py", font_size=16)
    set_shape_text(s.shapes[8], "Плюсы", font_size=20)
    set_shape_text(s.shapes[9], "Быстро, работает на любых фото без обучения.", font_size=16)

    # Слайд 7 — модель 3: глаза (MediaPipe)
    s = slides[6]
    set_shape_text(s.shapes[1], "Модель: глаза (MediaPipe)", font_size=34)
    set_shape_text(
        s.shapes[2],
        "MediaPipe Face Landmarker выдаёт лицевые ориентиры и blendshapes.\n"
        "Ключевые признаки: eyeBlink_L / eyeBlink_R (0=открыт, 1=закрыт).\n"
        "Решение: оба blink ≥ порога → «закрытые глаза».",
        font_size=16,
    )

    # Слайд 8 — модель 4: глаза (ResNet-18) + обучение
    s = slides[7]
    set_shape_text(s.shapes[1], "Модель: глаза (ResNet-18)", font_size=34)
    res = metrics.get("eye_state_resnet18", {})
    set_shape_text(
        s.shapes[4],
        "Архитектура: ResNet-18, последний слой на 2 класса (open/closed).\n"
        "Данные: MRL Eye Dataset → train/val.\n"
        "Обучение: src/training/train_eye_classifier.py → models/eye_state_resnet18.pth + CSV лог.\n"
        f"Итог (epoch {res.get('epoch', '-')}) val_acc={res.get('val_acc', '-')}, val_loss={res.get('val_loss', '-')}.",
        font_size=16,
    )
    set_shape_text(s.shapes[7], "По умолчанию в приложении ResNet выключен (лучше работает MediaPipe).", font_size=16)

    # Слайд 9 — графики обучения (вставка картинок)
    s = slides[8]
    set_shape_text(s.shapes[1], "Графики обучения двух моделей", font_size=34)
    set_shape_text(
        s.shapes[4],
        "ResNet-18: accuracy/loss по эпохам\n"
        "MobileNetV3: accuracy/loss по эпохам\n"
        "Графики строятся из CSV логов обучения автоматически.",
        font_size=14,
    )
    assets = ROOT / "assets"
    acc_png = assets / "resnet18_accuracy.png"
    loss_png = assets / "resnet18_loss.png"
    m_acc_png = assets / "mobilenetv3_accuracy.png"
    m_loss_png = assets / "mobilenetv3_loss.png"
    # Картинки добавляем поверх слайда (если есть); если нет — оставляем инструкцию.
    if acc_png.is_file():
        s.shapes.add_picture(str(acc_png), Inches(0.7), Inches(2.5), width=Inches(5.4))
    if loss_png.is_file():
        s.shapes.add_picture(str(loss_png), Inches(6.0), Inches(2.5), width=Inches(5.4))
    if m_acc_png.is_file():
        s.shapes.add_picture(str(m_acc_png), Inches(0.7), Inches(5.0), width=Inches(5.4))
    if m_loss_png.is_file():
        s.shapes.add_picture(str(m_loss_png), Inches(6.0), Inches(5.0), width=Inches(5.4))

    # Слайд 10 — метрики (по требованию)
    s = slides[9]
    set_shape_text(s.shapes[1], "Метрики качества (реальные, по обучению)", font_size=34)
    set_shape_text(s.shapes[2], "Основные метрики", font_size=20)
    mob = metrics.get("eye_state_mobilenetv3", {})
    set_shape_text(
        s.shapes[3],
        f"ResNet-18: train_acc={res.get('train_acc', '-')}, val_acc={res.get('val_acc', '-')}, "
        f"train_loss={res.get('train_loss', '-')}, val_loss={res.get('val_loss', '-')}.\n"
        f"MobileNetV3: train_acc={mob.get('train_acc', '-')}, val_acc={mob.get('val_acc', '-')}, "
        f"train_loss={mob.get('train_loss', '-')}, val_loss={mob.get('val_loss', '-')}.\n"
        "Вывод: сравниваем валидационную точность и устойчивость loss.",
        font_size=16,
    )

    # Слайд 11 — сравнение режимов (по требованию)
    s = slides[10]
    set_shape_text(s.shapes[1], "Сравнение моделей для open/closed eyes", font_size=34)
    # На шаблоне здесь текстовый заголовок; добавим поясняющий текстбокс.
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.0), Inches(4.5))
    set_shape_text(
        box,
        f"ResNet-18: val_acc={res.get('val_acc', '-')}, val_loss={res.get('val_loss', '-')}.\n"
        f"MobileNetV3: val_acc={mob.get('val_acc', '-')}, val_loss={mob.get('val_loss', '-')}.\n\n"
        "На этом прогоне ResNet-18 показал лучшую валидационную точность.\n"
        "В приложении основной режим глаз — MediaPipe; CNN остаётся дополнительной проверкой.",
        font_size=18,
    )
    summary_png = assets / "eye_models_summary.png"
    if summary_png.is_file():
        s.shapes.add_picture(str(summary_png), Inches(0.9), Inches(4.7), width=Inches(10.8))

    # Слайд 12 — «база данных» и проценты (по требованию)
    s = slides[11]
    set_shape_text(s.shapes[1], "Данные, проценты и «база»", font_size=34)
    set_shape_text(s.shapes[2], "Разбиение train/val", font_size=20)
    set_shape_text(s.shapes[3], f"Фактически: {train_pct:.1f}% / {val_pct:.1f}% (Train/Val).", font_size=16)
    set_shape_text(s.shapes[4], "Доли классов", font_size=20)
    set_shape_text(
        s.shapes[5],
        f"open = {open_total} ({open_pct:.1f}%)\n"
        f"closed = {closed_total} ({closed_pct:.1f}%)\n"
        "Статистика берётся из подготовленного датасета (папки train/val).",
        font_size=16,
    )
    set_shape_text(s.shapes[6], "Хранение результатов анализа", font_size=20)
    set_shape_text(
        s.shapes[7],
        "Файлы раскладываются по папкам.\n"
        "Это и есть простая «БД»: структура каталогов + причины.\n"
        "При желании можно добавить экспорт отчёта в CSV.",
        font_size=16,
    )

    # Слайд 13 — примеры причин
    s = slides[12]
    set_shape_text(s.shapes[1], "Примеры причин в «Неудачные»", font_size=34)

    # Слайд 14 — архитектура пайплайна (вместо старых метрик)
    s = slides[13]
    set_shape_text(s.shapes[1], "Пайплайн обработки", font_size=34)
    set_shape_text(s.shapes[3], "1) Глаза (MediaPipe)\n2) Экспозиция\n3) Резкость\n→ решение по кадру", font_size=18)
    set_shape_text(s.shapes[4], "Где настраивается", font_size=20)
    set_shape_text(s.shapes[5], "Раздел «Настройки» → app_settings.json", font_size=16)
    set_shape_text(s.shapes[6], "Где в коде", font_size=20)
    set_shape_text(s.shapes[7], "src/analysis/pipeline.py", font_size=16)

    # Слайд 15 — показ/демо
    s = slides[14]
    set_shape_text(s.shapes[1], "Показ (демонстрация)", font_size=34)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.0), Inches(3.5))
    set_shape_text(
        box,
        "1) Выбор папки с фото\n2) Запуск анализа\n3) Появление «Удачные» и «Неудачные/<причина>»\n"
        "4) Просмотр в разделе «Альбомы», ручной перенос при необходимости",
        font_size=20,
    )

    # Слайд 16 — заключение + ссылки
    s = slides[15]
    set_shape_text(s.shapes[1], "Заключение", font_size=34)
    set_shape_text(
        s.shapes[4],
        "Целевая аудитория: фотографы, фотостудии, мероприятия.\n"
        "Результат: ускорение первичного отбора кадров и единые критерии качества.\n"
        "Перспективы: новые критерии (улыбка/композиция), экспорт отчёта, оптимизация скорости.",
        font_size=18,
    )
    s.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.0), Inches(0.6)).text_frame.text = (
        "GitHub: https://github.com/Dorlian/Samsung_project"
    )

    pr.save(str(OUT_PPTX))
    print("Презентация:", OUT_PPTX)


def build_zapiska() -> None:
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Times New Roman"
    style.font.size = Pt(14)

    def add_center(lines: list[str]) -> None:
        for t in lines:
            p = doc.add_paragraph(t.strip())
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    def add_left(text: str) -> None:
        doc.add_paragraph(text)

    def add_heading(num: str, title: str) -> None:
        p = doc.add_paragraph()
        r = p.add_run(f"{num}. {title}")
        r.bold = True

    add_center(
        [
            "МИНИСТЕРСТВО НАУКИ И ВЫСШЕГО ОБРАЗОВАНИЯ РОССИЙСКОЙ ФЕДЕРАЦИИ",
            "ФЕДЕРАЛЬНОЕ ГОСУДАРСТВЕННОЕ БЮДЖЕТНОЕ ОБРАЗОВАТЕЛЬНОЕ УЧРЕЖДЕНИЕ ВЫСШЕГО ОБРАЗОВАНИЯ",
            "«МАГНИТОГОРСКИЙ ГОСУДАРСТВЕННЫЙ ТЕХНИЧЕСКИЙ",
            "УНИВЕРСИТЕТ ИМ. Г.И. НОСОВА»",
            "(ФГБОУ ВО «МГТУ ИМ. Г.И.НОСОВА»)",
            "Институт энергетики и автоматизированных систем",
            "Кафедра бизнес – информатики и информационных технологий",
        ]
    )
    doc.add_paragraph()
    p = doc.add_paragraph()
    r = p.add_run("ПРОЕКТНОЕ ЗАДАНИЕ")
    r.bold = True
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    add_center(
        [
            "по дисциплине: Компьютерное зрение и интеллектуальные системы",
            "на тему: Разработка десктоп-приложения «Фотоассистент фотографа» для автоматической сортировки "
            "фотографий по экспозиции, резкости и признаку закрытых глаз",
            "Исполнители: _________________________, студенты ___ курса, группа ___________",
            "Руководитель: _________________________, _________________________________",
            'Работа допущена к защите «___» _____________ 20__ г. __________',
            "(подпись)",
            'Работа защищена «___» _____________ 20__ г. с оценкой __________     __________',
            "    (оценка)         (подпись)",
            "Магнитогорск, 2026 г.",
        ]
    )

    doc.add_paragraph()
    add_left("Содержание")
    for line in (
        "Введение",
        "Описание структуры проекта",
        "Инструкция по установке и запуску",
        "Описание датасета",
        "Архитектура моделей и обучение",
        "Интеграция и пользовательский интерфейс",
        "Заключение",
        "Список источников",
    ):
        doc.add_paragraph(line)

    add_heading("1", "Введение")
    add_left(
        "Объём цифровой фотографии неуклонно растёт: съёмки мероприятий, портреты, репортажи сопровождаются "
        "сотнями и тысячами кадров. Значительная часть кадров не пригодна для передачи заказчику из‑за технических "
        "дефектов: неудачная экспозиция, смаз, закрытые глаза на портретах. Ручной просмотр каждого файла "
        "отнимает много времени и приводит к утомлению оператора."
        "\n\n"
        "Методы компьютерного зрения и машинного обучения позволяют автоматизировать первичный отбор: оценить "
        "распределение яркости, резкость изображения, наличие лица и состояние век по данным модели лицевых ориентиров."
        "\n\n"
        "Цель проекта — разработка десктоп-приложения на языке Python с графическим интерфейсом, выполняющего "
        "рекурсивный обход выбранной папки с файлами .jpg / .jpeg / .png и распределение снимков по каталогам "
        "«Удачные» и «Неудачные» с указанием причины отбраковки."
        "\n\n"
        "Задачи: анализ предметной области; выбор и настройка метрик; подготовка данных и обучение вспомогательного "
        "классификатора глаз (опционально); интеграция MediaPipe Face Landmarker; разработка интерфейса и тестирование."
    )

    add_heading("2", "Описание структуры проекта")
    add_left(
        "Исходный код и документация размещены в открытом репозитории GitHub: https://github.com/Dorlian/Samsung_project"
        "\n\n"
        "Корень проекта: README.md и .bat; app/__main__.py (python -m app), requirements/requirements.txt, "
        "bin/run.sh, build_exe.bat и packaging/photo_assistant.spec (PyInstaller), каталог models/ для MediaPipe "
        "и опционального файла весов eye_state_resnet18.pth, каталог training/ с инструкциями по данным для обучения."
        "\n\n"
        "Пакет src/analysis содержит модули экспозиции, резкости, анализа глаз и сборку пайплайна PhotoAnalyzer; "
        "src/config — сохранение порогов в app_settings.json; src/models — архитектура ResNet для глаз; "
        "src/training — скрипты prepare_mrl_dataset.py и train_eye_classifier.py; src/ui — раздел «Альбомы»; "
        "src/utils — обход файлов и перемещение в папки результатов."
    )

    add_heading("3", "Инструкция по установке и запуску")
    add_left(
        "Требуется Python 3.10–3.12. Рекомендуется создание виртуального окружения и установка зависимостей командой "
        "pip install -r requirements/requirements.txt. Запуск: python -m app или run.bat / bin/run.sh (см. README.md)."
        "\n\n"
        "При первом анализе изображений с лицами в каталог models/ загружается файл face_landmarker.task (необходим "
        "доступ в Интернет один раз). Опциональная нейросеть для глаз по умолчанию отключена в настройках; при "
        "включении и наличии файла весов используется дополнительная проверка кропов глаз."
    )

    add_heading("4", "Описание датасета")
    add_left(
        "Для обучения классификатора «открытый / закрытый глаз» применяется открытый MRL Eye Dataset "
        "(изображения области глаза с метками в имени файла). Официальный источник: https://mrl.cs.vsb.cz/eyedataset. "
        "Скрипт prepare_mrl_dataset.py формирует структуры train/open, train/closed, val/open, val/closed. "
        "Полный объём датасета в репозиторий не включён из‑за размера; инструкции по скачиванию приведены в training/README.md."
        "\n\n"
        "Критерии экспозиции и резкости не требуют отдельного датасета: используются классические вычисления по пикселям изображения."
    )

    add_heading("5", "Архитектура моделей и обучение")
    add_left(
        "Детекция лица и оценка смыкания век выполняется средствами MediaPipe Face Landmarker с использованием "
        "blendshape-категорий eyeBlink_L и eyeBlink_R."
        "\n\n"
        "Дополнительно реализован классификатор на базе ResNet-18 (два выхода). Обучение — скрипт train_eye_classifier.py "
        "с параметрами по умолчанию: размер входа 64×64, пакет 64, до 10 эпох, оптимизатор Adam, функция потерь — "
        "перекрёстная энтропия; лучшие веса сохраняются в models/eye_state_resnet18.pth."
        "\n\n"
        "Метрики экспозиции и резкости задаются порогами в пользовательском интерфейсе и не используют обучение на размеченном датасете."
    )

    add_heading("6", "Интеграция и пользовательский интерфейс")
    add_left(
        "Интерфейс реализован на библиотеке Tkinter. Раздел «Анализ»: выбор папки, запуск и остановка обработки, "
        "прогресс и журнал. Раздел «Альбомы»: просмотр миниатюр и ручной перенос выбранных файлов между альбомами. "
        "Раздел «Настройки»: пороги экспозиции и резкости, включение опциональной CNN для глаз; параметры сохраняются в app_settings.json."
        "\n\n"
        "Модуль pipeline.py последовательно вызывает проверку глаз, экспозиции и резкости и формирует итоговое решение по каждому файлу."
    )

    add_heading("7", "Заключение")
    add_left(
        "Разработано приложение для автоматической сортировки фотографий по экспозиции, резкости и признаку закрытых глаз. "
        "Реализованы классические методы анализа изображения, детекция лица MediaPipe и опциональное обучение свёрточного классификатора."
        "\n\n"
        "Практическая ценность — сокращение времени первичного отбора кадров. Перспективы: добавление новых критериев, "
        "улучшение устойчивости к освещению и ракурсу, оптимизация производительности."
    )

    add_heading("8", "Список использованных источников")
    for i, src in enumerate(
        [
            "MediaPipe Face Landmarker — https://ai.google.dev/edge/mediapipe/solutions/vision/face_landmarker",
            "MRL Eye Dataset — https://mrl.cs.vsb.cz/eyedataset",
            "Документация PyTorch — https://pytorch.org/docs/stable/index.html",
            "Документация OpenCV — https://docs.opencv.org/",
            "Документация Python Tkinter — https://docs.python.org/3/library/tkinter.html",
        ],
        start=1,
    ):
        doc.add_paragraph(f"{i}. {src}")

    doc.save(str(OUT_DOCX))
    print("Записка:", OUT_DOCX)


def main() -> None:
    build_zapiska()
    build_presentation()


if __name__ == "__main__":
    main()
