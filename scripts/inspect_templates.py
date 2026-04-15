"""Dump template pptx / docx text to UTF-8 files for authoring."""
from __future__ import annotations

import re
import zipfile
from pathlib import Path

from pptx import Presentation


def dump_pptx(path: Path, out: Path) -> None:
    pr = Presentation(str(path))
    lines: list[str] = []
    for i, slide in enumerate(pr.slides):
        lines.append(f"=== SLIDE {i + 1} ===")
        for j, sh in enumerate(slide.shapes):
            if not hasattr(sh, "text"):
                continue
            t = (sh.text or "").strip()
            if not t:
                continue
            lines.append(f"  shape[{j}] ({sh.shape_type}):\n{t}")
        lines.append("")
    out.write_text("\n".join(lines), encoding="utf-8")


def dump_docx(path: Path, out: Path) -> None:
    z = zipfile.ZipFile(path)
    xml = z.read("word/document.xml").decode("utf-8")
    text = re.sub(r"</w:p>", "\n", xml)
    text = re.sub(r"<[^>]+>", "", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    out.write_text(text.strip() + "\n", encoding="utf-8")


if __name__ == "__main__":
    root = Path(__file__).resolve().parents[1]
    dump_pptx(
        Path(r"c:\Users\studi\Downloads\presentation-template.pptx"),
        root / "_inspect_pptx.txt",
    )
    dump_docx(
        Path(r"c:\Users\studi\Downloads\zapiska (1).docx"),
        root / "_inspect_zapiska1.txt",
    )
    print("ok", root / "_inspect_pptx.txt", root / "_inspect_zapiska1.txt")
